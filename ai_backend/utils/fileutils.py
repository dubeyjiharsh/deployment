# Standard library imports
import os
import json
import io
import logging
from datetime import datetime
from urllib.parse import urlparse, unquote
from typing import Dict, List, Tuple, Optional, Union, Any
from dotenv import load_dotenv
import pandas as pd
from openai import AzureOpenAI

# File Processing imports
import pdfplumber
from pptx import Presentation
from docx import Document

# Azure SDK imports
from azure.storage.blob import BlobServiceClient 
from azure.data.tables import TableServiceClient
from azure.core.exceptions import ResourceExistsError
from azure.core.exceptions import AzureError
from azure.core.credentials import AzureNamedKeyCredential

# Custom logging utility functions.
from utils.logutils import get_logger, log_context  

# Initialize environment variables
load_dotenv()

# Get logger from logutils
logger = get_logger()

# Proxy Configurations (Optional / Set here if required)
# os.environ["HTTP_PROXY"] = os.getenv("HTTP_PROXY")
# os.environ["HTTPS_PROXY"] = os.getenv("HTTPS_PROXY")
# os.environ["SSL_CERT_FILE"] = os.getenv("SSL_CERT_FILE")
# os.environ["NO_PROXY"] = os.getenv("NO_PROXY")

# Azure OpenAI Configuration
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION")
AZURE_OPENAI_MODEL = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME")
AZURE_OPENAI_TEMPERATURE = float(os.getenv("AZURE_OPENAI_TEMPERATURE", "0.6"))  # Default to 0.6 if not specified

# Azure Blob Storage Configuration
AZURE_BLOB_STORAGE_NAME = os.getenv("AZURE_BLOB_STORAGE_NAME")
AZURE_BLOB_STORAGE_KEY = os.getenv("AZURE_BLOB_STORAGE_KEY")
AZURE_BLOB_STORAGE_ENDPOINT = os.getenv("AZURE_BLOB_STORAGE_ENDPOINT")

# Document Processing Control Configuration
MAX_SUMMARY_TOKENS = int(os.getenv("MAX_SUMMARY_TOKENS", "200"))  # Max tokens in generated summary
MAX_PDF_PAGES = int(os.getenv("MAX_PDF_PAGES", "20"))  # Max pages to process in PDF
MAX_WORD_PAGES = int(os.getenv("MAX_WORD_PAGES", "20"))  # Max paragraphs to process in Word
MAX_EXCEL_SHEETS = int(os.getenv("MAX_EXCEL_SHEETS", "10"))  # Max sheets to process in Excel
MAX_PPT_SLIDES = int(os.getenv("MAX_PPT_SLIDES", "10"))  # Max slides to process in PowerPoint
MAX_CONTENT_TOKENS = int(os.getenv("MAX_CONTENT_TOKENS", "2000"))  # Max tokens for content to summarize

# Supported file formats
ALLOWED_FORMATS = [".pdf", ".ppt", ".pptx", ".doc", ".docx", ".xls", ".xlsx"]

def get_blob_service_client() -> BlobServiceClient:
    """
    Creates and returns an Azure Blob Service Client.
    
    Returns:
        BlobServiceClient: Initialized Azure Blob Storage client
        
    Raises:
        ValueError: If Azure Blob Storage credentials are missing
    """
    if not all([AZURE_BLOB_STORAGE_ENDPOINT, AZURE_BLOB_STORAGE_KEY]):
        error_msg = "Azure Blob Storage credentials are missing in environment variables"
        log_context('error', 'system', 'system', error_msg)
        raise ValueError(error_msg)
        
    try:
        return BlobServiceClient(
            account_url=AZURE_BLOB_STORAGE_ENDPOINT,
            credential=AZURE_BLOB_STORAGE_KEY
        )
    except Exception as e:
        log_context('error', 'system', 'system', f"Failed to create blob service client: {str(e)}")
        raise

def save_files_to_blob_storage(user_name: str, session_id: str, files: List[Tuple[str, bytes]]) -> List[str]:
    """
    Save multiple files to Azure Blob Storage using a virtual directory hierarchy.
    
    Args:
        user_name (str): Name of the root container (usually user identifier)
        session_id (str): Unique identifier for the virtual directory
        files (list): List of tuples containing (filename, file_contents)
        
    Returns:
        list: List of full URLs of successfully saved files in blob storage
        
    Raises:
        ValueError: If credentials are invalid or parameters are incorrect
        Exception: For any other errors during file upload process
    """
    saved_files = []
    
    try:
        # Initialize blob service client
        blob_service_client = get_blob_service_client()
        
        # Get or create user container
        container_client = blob_service_client.get_container_client(container=user_name)
        if not container_client.exists():
            container_client.create_container()
            log_context('info', user_name, session_id, "Created new container for user")
        else:
            log_context('info', user_name, session_id, "Using existing container for user")
        
        # Process each file
        for file_name, file_content in files:
            # Validate file extension
            file_extension = os.path.splitext(file_name)[1].lower()
            if file_extension not in ALLOWED_FORMATS:
                log_context('warning', user_name, session_id, f"Unsupported file format: {file_name}")
                continue
            
            try:
                # Create virtual path for the blob
                blob_path = f"{session_id}/{file_name}"
                
                # Upload file content
                blob_client = container_client.get_blob_client(blob=blob_path)
                blob_client.upload_blob(file_content, overwrite=True)
                
                # Construct the complete URL
                complete_url = f"{AZURE_BLOB_STORAGE_ENDPOINT}{user_name}/{blob_path}"
                
                saved_files.append(complete_url)
                log_context('info', user_name, session_id, f"Successfully uploaded: {file_name}")

            except Exception as e:
                log_context('error', user_name, session_id, f"Error uploading {file_name}: {str(e)}")
                # Continue with other files even if one fails
                
        return saved_files
        
    except ValueError as ve:
        log_context('error', user_name, session_id, f"Configuration error: {str(ve)}")
        raise
    except Exception as e:
        log_context('error', user_name, session_id, f"Blob storage error: {str(e)}")
        raise

def get_blob_data(blob_url: str, user_name: str = "system", session_id: str = "blob-operation") -> bytes:
    """
    Download file content from Azure Blob Storage.
    
    Args:
        blob_url (str): Complete URL to the blob
        user_name (str, optional): User name for logging context. Defaults to "system".
        session_id (str, optional): Session ID for logging context. Defaults to "blob-operation".
        
    Returns:
        bytes: File content
        
    Raises:
        ValueError: If the URL is invalid
        AzureError: If there's an issue with the Azure Blob service
        Exception: For any other errors during download
    """
    if not blob_url:
        error_msg = "Blob URL cannot be empty"
        log_context('error', user_name, session_id, error_msg)
        raise ValueError(error_msg)
        
    try:
        # Parse the URL to extract container and blob path
        parsed_url = urlparse(blob_url)
        path_parts = parsed_url.path.strip('/').split('/')
        
        if len(path_parts) < 2:
            error_msg = f"Invalid blob URL format: {blob_url}"
            log_context('error', user_name, session_id, error_msg)
            raise ValueError(error_msg)
            
        container_name = path_parts[0]
        blob_path = '/'.join(path_parts[1:])
        blob_path = unquote(blob_path)  # Handle URL encoded characters
        
        # Get container and blob client
        blob_service_client = get_blob_service_client()
        container_client = blob_service_client.get_container_client(container_name)
        blob_client = container_client.get_blob_client(blob_path)
        
        # Download blob content
        blob_data = blob_client.download_blob().read()
        log_context('info', user_name, session_id, f"Successfully downloaded blob from: {blob_url}")
        return blob_data
        
    except AzureError as e:
        error_msg = f"Azure Blob Storage error: {str(e)}"
        log_context('error', user_name, session_id, error_msg)
        raise
    except Exception as e:
        error_msg = f"Error downloading blob: {str(e)}"
        log_context('error', user_name, session_id, error_msg)
        raise

def get_azure_openai_client() -> AzureOpenAI:
    """
    Initialize and return Azure OpenAI client.
    
    Returns:
        AzureOpenAI: Initialized Azure OpenAI client
        
    Raises:
        ValueError: If OpenAI credentials are missing
        Exception: For any other initialization errors
    """
    if not all([AZURE_OPENAI_ENDPOINT, AZURE_OPENAI_API_KEY, AZURE_OPENAI_API_VERSION]):
        error_msg = "Azure OpenAI credentials are incomplete in environment variables"
        log_context('error', 'system', 'system', error_msg)
        raise ValueError(error_msg)
        
    try:
        return AzureOpenAI(
            azure_endpoint=AZURE_OPENAI_ENDPOINT,
            api_key=AZURE_OPENAI_API_KEY,
            api_version=AZURE_OPENAI_API_VERSION
        )
    except Exception as e:
        log_context('error', 'system', 'system', f"Error initializing Azure OpenAI client: {str(e)}")
        raise


def truncate_content_for_summary(content: str, max_tokens: int = MAX_CONTENT_TOKENS) -> str:
    """
    Truncate content to a maximum number of tokens (approximate).
    
    Args:
        content (str): The content to truncate
        max_tokens (int, optional): Maximum number of tokens. Defaults to MAX_CONTENT_TOKENS.
        
    Returns:
        str: Truncated content
    """
    # A very rough estimate: 1 token ≈ 4 characters in English
    # This is an approximation - actual tokenization varies by model
    estimated_char_limit = max_tokens * 4
    
    if len(content) <= estimated_char_limit:
        return content
        
    logger.info(f"Content truncated from {len(content)} chars to approximately {estimated_char_limit} chars")
    return content[:estimated_char_limit] + "... [content truncated due to length]"


def generate_summary(content: str, max_tokens: int = MAX_SUMMARY_TOKENS, model: Optional[str] = None,
                user_name: str = "system", session_id: str = "system") -> str:
    """
    Generate a concise summary of content using Azure OpenAI.
    
    Args:
        content (str): The text content to summarize
        max_tokens (int, optional): Maximum length of the summary. Defaults to MAX_SUMMARY_TOKENS.
        model (str, optional): Model to use. Defaults to model from env vars.
        user_name (str, optional): User identifier for logging. Defaults to "system".
        session_id (str, optional): Session identifier for logging. Defaults to "default".
    
    Returns:
        str: Generated summary
        
    Raises:
        ValueError: If content is empty or model is not specified
        Exception: For any other errors during summary generation
    """
    if not content:
        error_msg = "Cannot generate summary for empty content"
        log_context('error', user_name, session_id, error_msg)
        raise ValueError(error_msg)
        
    if not model:
        model = AZURE_OPENAI_MODEL
        
    if not model:
        error_msg = "No model specified for summary generation"
        log_context('error', user_name, session_id, error_msg)
        raise ValueError(error_msg)
    
    # Truncate content if needed
    content = truncate_content_for_summary(content)
    
    try:
        azure_client = get_azure_openai_client()
        log_context('info', user_name, session_id, f"Generating summary using model: {model} with max_tokens: {max_tokens}")
        
        response = azure_client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "You are an assistant providing concise summaries."},
                {"role": "user", "content": f"Provide a concise summary focusing on key insights:\n\n{content}"}
            ],
            max_tokens=max_tokens,
            temperature=AZURE_OPENAI_TEMPERATURE
        )
        return response.choices[0].message.content
    except Exception as e:
        error_msg = f"Error generating summary: {str(e)}"
        log_context('error', user_name, session_id, error_msg)
        raise

#----------------------------------------------------------------------------------

def analyze_word_file(doc_path: str, max_tokens: int = MAX_SUMMARY_TOKENS,
                     user_name: str = "system", session_id: str = "system") -> Dict[str, str]:
    """
    Analyze Word documents and generate summaries.
    
    Args:
        doc_path (str): Path or URL to the Word document
        max_tokens (int, optional): Maximum length of summary. Defaults to MAX_SUMMARY_TOKENS.
        user_name (str, optional): User identifier for logging. Defaults to "system".
        session_id (str, optional): Session identifier for logging. Defaults to "default".
    
    Returns:
        Dict[str, str]: Dictionary with filename as key and summary as value
    """
    file_name = doc_path.split('/')[-1]
    log_context('info', user_name, session_id, f"Starting analysis of Word file: {file_name}")
    
    try:
        doc_data = get_blob_data(doc_path, user_name, session_id)
        doc_buffer = io.BytesIO(doc_data)
        content = ""

        # Process based on file extension - Note: python-docx operations are synchronous
        if file_name.lower().endswith('.docx'):
            doc = Document(doc_buffer)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            # Limit number of paragraphs to process
            if len(paragraphs) > MAX_WORD_PAGES:
                log_context('info', user_name, session_id, 
                            f"Limiting Word document analysis to first {MAX_WORD_PAGES} of {len(paragraphs)} paragraphs")
                paragraphs = paragraphs[:MAX_WORD_PAGES]
                
            content = "\n".join(paragraphs)
        else:
            # For older .doc format, this requires textract
            # If textract is not available, provide an error message
            error_msg = "Textract module is required for .doc files but is not available"
            log_context('error', user_name, session_id, error_msg)
            return {file_name: "Processing .doc files requires the textract library"}

        summary = generate_summary(content, max_tokens, AZURE_OPENAI_MODEL, user_name, session_id)
        return {file_name: summary}
    
    except Exception as e:
        error_msg = f"Unexpected error during Word file analysis: {str(e)}"
        log_context('error', user_name, session_id, error_msg)
        return {file_name: f"Error analyzing file: {str(e)}"}
    
#----------------------------------------------------------------------------------

def analyze_excel_file(excel_path: str, max_tokens: int = MAX_SUMMARY_TOKENS,
                  user_name: str = "system", session_id: str = "system") -> Dict[str, str]:
    """
    Analyze Excel files and generate summaries with proper table formatting.
    
    Args:
        excel_path (str): Path or URL to the Excel file
        max_tokens (int, optional): Maximum length of summary. Defaults to MAX_SUMMARY_TOKENS.
        user_name (str, optional): User identifier for logging. Defaults to "system".
        session_id (str, optional): Session identifier for logging. Defaults to "default".
    
    Returns:
        Dict[str, str]: Dictionary with filename as key and summary as value
    """
    file_name = excel_path.split('/')[-1]
    log_context('info', user_name, session_id, f"Starting analysis of Excel file: {file_name}")
    
    try:
        # Download and read the Excel file
        excel_data = get_blob_data(excel_path, user_name, session_id)
        excel_buffer = io.BytesIO(excel_data)
        
        with pd.ExcelFile(excel_buffer, engine='openpyxl') as xlsx:
            sheet_info = []
            
            # Limit the number of sheets to process
            sheet_names = xlsx.sheet_names[:MAX_EXCEL_SHEETS]
            if len(xlsx.sheet_names) > MAX_EXCEL_SHEETS:
                log_context('info', user_name, session_id,
                         f"Limiting analysis to first {MAX_EXCEL_SHEETS} of {len(xlsx.sheet_names)} sheets")
            
            for sheet_name in sheet_names:
                try:
                    df = pd.read_excel(xlsx, sheet_name=sheet_name, na_filter=True)
                    if not df.empty:
                        # Format the data as a markdown table
                        md_table = format_dataframe_as_table(df)
                        
                        sheet_info.append(
                            f"## Sheet: {sheet_name}\n"
                            f"Rows: {len(df)}, Columns: {len(df.columns)}\n\n"
                            f"{md_table}\n"
                        )
                except Exception as e:
                    error_msg = f"Error processing sheet {sheet_name}: {str(e)}"
                    log_context('error', user_name, session_id, error_msg)
                    sheet_info.append(f"Sheet: {sheet_name} (Error: {str(e)})\n")
            
            content = "\n".join(sheet_info)
            summary = generate_summary(content, max_tokens, AZURE_OPENAI_MODEL, user_name, session_id)
            
            return {file_name: summary}
    
    except Exception as e:
        error_msg = f"Unexpected error during Excel analysis: {str(e)}"
        log_context('error', user_name, session_id, error_msg)
        return {file_name: f"Error analyzing file: {str(e)}"}

def format_dataframe_as_table(df, max_rows=10):
    """
    Format a pandas DataFrame as a markdown table.
    
    Args:
        df (pd.DataFrame): The DataFrame to format
        max_rows (int): Maximum number of rows to include
    
    Returns:
        str: Markdown formatted table
    """
    # Limit the number of rows to prevent oversized tables
    if len(df) > max_rows:
        df_display = df.head(max_rows)
        footer = f"*Showing {max_rows} of {len(df)} rows*"
    else:
        df_display = df
        footer = ""
    
    # Handle potential display issues with large columns
    max_col_length = 50
    for col in df_display.columns:
        if df_display[col].dtype == 'object':  # String columns
            df_display[col] = df_display[col].astype(str).apply(
                lambda x: (x[:max_col_length] + '...') if len(x) > max_col_length else x
            )
    
    # Create markdown table header
    header = "| " + " | ".join(str(col) for col in df_display.columns) + " |"
    separator = "| " + " | ".join(["---"] * len(df_display.columns)) + " |"
    
    # Create table rows
    rows = []
    for _, row in df_display.iterrows():
        # Format values to handle different data types appropriately
        formatted_values = []
        for val in row:
            if pd.isna(val):
                formatted_values.append("")
            elif isinstance(val, (int, float)):
                formatted_values.append(str(val))
            else:
                formatted_values.append(str(val).replace("|", "\\|"))  # Escape pipe characters
        
        rows.append("| " + " | ".join(formatted_values) + " |")
    
    # Combine all parts of the table
    table = "\n".join([header, separator] + rows)
    
    if footer:
        table += f"\n\n{footer}"
    
    return table

#----------------------------------------------------------------------------------

def analyze_powerpoint_file(ppt_path: str, max_tokens: int = MAX_SUMMARY_TOKENS,
                      user_name: str = "system", session_id: str = "system") -> Dict[str, str]:
    """
    Analyze PowerPoint presentations and generate summaries with enhanced error handling.
    
    Args:
        ppt_path (str): Path or URL to the PowerPoint file
        max_tokens (int, optional): Maximum length of summary. Defaults to MAX_SUMMARY_TOKENS.
        user_name (str, optional): User identifier for logging. Defaults to "system".
        session_id (str, optional): Session identifier for logging. Defaults to "default".
    
    Returns:
        Dict[str, str]: Dictionary with filename as key and summary as value
    """
    file_name = ppt_path.split('/')[-1]
    log_context('info', user_name, session_id, f"Starting analysis of PowerPoint file: {file_name}")
    
    try:
        # Get file data from blob storage
        ppt_data = get_blob_data(ppt_path, user_name, session_id)
        ppt_buffer = io.BytesIO(ppt_data)
        
        # Try the primary method first, with fallback options if it fails
        try:
            content = extract_ppt_content_primary_method(ppt_buffer, user_name, session_id)
        except Exception as primary_error:
            # Log the primary method failure
            log_context('warning', user_name, session_id, 
                       f"Primary PowerPoint extraction method failed: {str(primary_error)}")
            
            # Reset buffer position for next attempt
            ppt_buffer.seek(0)
            
            try:
                # Try fallback method
                content = extract_ppt_content_fallback_method(ppt_buffer, user_name, session_id)
            except Exception as fallback_error:
                # Both methods failed
                error_msg = f"All PowerPoint extraction methods failed. Primary error: {str(primary_error)}. Fallback error: {str(fallback_error)}"
                log_context('error', user_name, session_id, error_msg)
                return {file_name: f"Unable to extract content from PowerPoint: {str(fallback_error)}"}
        
        # Generate summary if we have content
        if content:
            full_content = "\n\n".join(content)
            summary = generate_summary(full_content, max_tokens, AZURE_OPENAI_MODEL, user_name, session_id)
            return {file_name: summary}
        else:
            return {file_name: "No extractable text found in presentation"}
    
    except Exception as e:
        # Catch-all for any other errors
        error_msg = f"Unexpected error during PowerPoint file analysis: {str(e)}"
        log_context('error', user_name, session_id, error_msg)
        return {file_name: f"Error analyzing file: {str(e)}"}


def extract_ppt_content_primary_method(ppt_buffer, user_name: str, session_id: str):
    """
    Primary method to extract content from PowerPoint using the standard approach.
    
    Args:
        ppt_buffer: BytesIO buffer containing the PowerPoint file
        user_name: User identifier for logging
        session_id: Session identifier for logging
        
    Returns:
        list: List of strings with content from each slide
    """
    content = []
    
    try:
        # Create fresh buffer to ensure we're at the start
        ppt_buffer.seek(0)
        presentation = Presentation(ppt_buffer)
        
        # Limit number of slides to process
        max_slides = min(len(presentation.slides), MAX_PPT_SLIDES)
        
        # Process each slide with added safety
        for i, slide in enumerate(presentation.slides[:max_slides]):
            try:
                slide_text = []
                
                # IMPORTANT: Use safer iteration that doesn't rely on slide.shapes being iterable
                if hasattr(slide, "shapes"):
                    # Get the total shape count first to avoid iteration issues
                    shape_count = len(slide.shapes)
                    
                    # Loop through shapes by index to avoid iteration errors
                    for shape_idx in range(shape_count):
                        try:
                            shape = slide.shapes[shape_idx]
                            
                            # First method: extract from plain text attribute
                            if hasattr(shape, "text") and shape.text is not None:
                                text = shape.text.strip()
                                if text:
                                    slide_text.append(text)
                                continue
                                
                            # Second method: extract from text frames
                            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                                for paragraph in shape.text_frame.paragraphs:
                                    if paragraph.text and paragraph.text.strip():
                                        slide_text.append(paragraph.text.strip())
                                        
                            # Third method: handle tables
                            if hasattr(shape, "has_table") and shape.has_table:
                                if hasattr(shape, "table"):
                                    table_text = []
                                    for row in shape.table.rows:
                                        for cell in row.cells:
                                            if cell.text and cell.text.strip():
                                                table_text.append(cell.text.strip())
                                    if table_text:
                                        slide_text.append(" | ".join(table_text))
                                        
                        except Exception as shape_error:
                            # Log but continue with next shape
                            log_context('debug', user_name, session_id, 
                                       f"Error processing shape {shape_idx} in slide {i+1}: {str(shape_error)}")
                
                if slide_text:
                    content.append(f"Slide {i+1}: {' '.join(slide_text)}")
            except Exception as slide_error:
                # Log but continue with next slide
                log_context('warning', user_name, session_id, 
                           f"Error processing slide {i+1}: {str(slide_error)}")
                
        return content
    except Exception as e:
        # Let the caller handle this exception
        log_context('error', user_name, session_id, f"Primary extraction method failed: {str(e)}")
        raise


def extract_ppt_content_fallback_method(ppt_buffer, user_name: str, session_id: str):
    """
    Alternative method to extract content from PowerPoint with minimal dependencies on structure.
    This method uses a more conservative approach that's less likely to encounter errors.
    
    Args:
        ppt_buffer: BytesIO buffer containing the PowerPoint file
        user_name: User identifier for logging
        session_id: Session identifier for logging
        
    Returns:
        list: List of strings with content from each slide
    """
    from pptx import Presentation
    import xml.etree.ElementTree as ET
    from pptx.oxml import parse_xml
    
    content = []
    
    try:
        # Reset buffer position
        ppt_buffer.seek(0)
        
        # Load presentation package
        presentation = Presentation(ppt_buffer)
        
        # Get maximum slides to process
        try:
            slide_count = len(presentation.slides)
            max_slides = min(slide_count, MAX_PPT_SLIDES)
        except Exception:
            # If we can't determine slide count, use a conservative estimate
            max_slides = MAX_PPT_SLIDES
            log_context('warning', user_name, session_id, 
                       f"Could not determine slide count, using maximum: {MAX_PPT_SLIDES}")
        
        try:
            # Try to enumerate slides directly
            for i, slide in enumerate(presentation.slides[:max_slides]):
                slide_text = []
                
                # Use the presentation's XML directly
                try:
                    if hasattr(slide, "_element"):
                        # Find all text elements in the slide's XML
                        for element in slide._element.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                            if element.text and element.text.strip():
                                slide_text.append(element.text.strip())
                except Exception as xml_error:
                    log_context('warning', user_name, session_id, 
                               f"Error parsing slide XML: {str(xml_error)}")
                
                if slide_text:
                    content.append(f"Slide {i+1}: {' '.join(slide_text)}")
                else:
                    content.append(f"Slide {i+1}: [No extractable text found]")
                    
        except Exception as slide_enum_error:
            # If direct slide enumeration fails, try an even more basic XML approach
            log_context('warning', user_name, session_id, 
                       f"Slide enumeration failed: {str(slide_enum_error)}. Attempting XML-only extraction.")
            
            # Reset presentation for XML-only extraction
            ppt_buffer.seek(0)
            pres = Presentation(ppt_buffer)
            
            # Try to extract text from the entire presentation XML as a fallback
            if hasattr(pres, "part") and hasattr(pres.part, "package"):
                slide_parts = []
                
                # Try to get slide parts from package
                try:
                    for rel in pres.part.rels.values():
                        if rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide':
                            slide_parts.append(rel.target_part)
                except Exception:
                    log_context('warning', user_name, session_id, "Could not enumerate slides via relationships")
                
                # Process slide parts we found
                for i, slide_part in enumerate(slide_parts[:MAX_PPT_SLIDES]):
                    slide_text = []
                    
                    try:
                        # Extract text directly from slide's XML
                        slide_xml = slide_part.blob
                        if slide_xml:
                            try:
                                root = ET.fromstring(slide_xml)
                                # Look for text elements in the XML
                                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                                for t in root.findall('.//a:t', ns):
                                    if t.text and t.text.strip():
                                        slide_text.append(t.text.strip())
                            except ET.ParseError:
                                log_context('warning', user_name, session_id, 
                                           f"Error parsing XML for slide {i+1}")
                    except Exception as part_error:
                        log_context('warning', user_name, session_id, 
                                   f"Error processing slide part {i+1}: {str(part_error)}")
                        
                    if slide_text:
                        content.append(f"Slide {i+1}: {' '.join(slide_text)}")
                    
        return content
                
    except Exception as e:
        log_context('error', user_name, session_id, f"Fallback extraction method failed: {str(e)}")
        raise


#----------------------------------------------------------------------------------


def analyze_pdf_file(pdf_path: str, max_tokens: int = MAX_SUMMARY_TOKENS,
                      user_name: str = "system", session_id: str = "system") -> Dict[str, str]:
    """
    Analyze PDF files and generate summaries.
    
    Args:
        pdf_path (str): Path or URL to the PDF file
        max_tokens (int, optional): Maximum length of summary. Defaults to MAX_SUMMARY_TOKENS.
    
    Returns:
        Dict[str, str]: Dictionary with filename as key and summary as value
    """
    log_context('info', user_name, session_id, f"Starting analysis of PDF file: {pdf_path}")
    file_name = pdf_path.split('/')[-1]
    
    try:
        pdf_data = get_blob_data(pdf_path, user_name, session_id)
        pdf_buffer = io.BytesIO(pdf_data)
        
        # Extract text from PDF pages
        full_text = ""
        with pdfplumber.open(pdf_buffer) as pdf:
            # Get total pages for logging
            total_pages = len(pdf.pages)
            log_context('info', user_name, session_id, f"PDF has {total_pages} pages")
            
            # Limit the number of pages to process
            pages_to_process = min(total_pages, MAX_PDF_PAGES)
            
            if total_pages > MAX_PDF_PAGES:
                logger.info(f"Limiting analysis to first {MAX_PDF_PAGES} of {total_pages} pages")
                
            # Extract content from each page
            for i in range(pages_to_process):
                page = pdf.pages[i]
                text = page.extract_text()
                if text:
                    full_text += f"--- PAGE {i+1} ---\n{text.strip()}\n\n"
        
        # Only attempt summary if we extracted some text
        if full_text:
            summary = generate_summary(full_text, max_tokens, AZURE_OPENAI_MODEL, user_name, session_id)
            return {file_name: summary}
        else:
            error_msg = "No extractable text found in PDF"
            log_context('info', user_name, session_id, f"No extractable text found in PDF: {file_name}")
            return {file_name: error_msg}

    except Exception as e:
        error_msg = f"Unexpected error during PDF file analysis: {str(e)}"
        log_context('error', user_name, session_id, error_msg)
        return {file_name: f"Error analyzing file: {str(e)}"}


def process_document(document_url: str, max_tokens: int = MAX_SUMMARY_TOKENS, 
                 user_name: str = "system", session_id: str = "system") -> Dict[str, str]:
    """
    Process a document based on its file extension.
    
    Args:
        document_url (str): URL to the document in blob storage
        max_tokens (int, optional): Maximum length of summary. Defaults to MAX_SUMMARY_TOKENS.
        user_name (str, optional): User identifier for logging. Defaults to "system".
        session_id (str, optional): Session identifier for logging. Defaults to "default".
        
    Returns:
        Dict[str, str]: Dictionary with filename as key and summary as value
    """
    file_name = document_url.split('/')[-1]
    file_extension = os.path.splitext(file_name)[1].lower()
    
    try:
        if file_extension in ['.xls', '.xlsx']:
            return analyze_excel_file(document_url, max_tokens, user_name, session_id)
        elif file_extension in ['.doc', '.docx']:
            return analyze_word_file(document_url, max_tokens, user_name, session_id)
        elif file_extension in ['.ppt', '.pptx']:
            return analyze_powerpoint_file(document_url, max_tokens, user_name, session_id)
        elif file_extension == '.pdf':
            return analyze_pdf_file(document_url, max_tokens, user_name, session_id)
        elif file_extension == '.txt':
            # For text files, download and analyze directly
            try:
                log_context('info', user_name, session_id, f"Processing text file: {file_name}")
                text_data = get_blob_data(document_url, user_name, session_id).decode('utf-8', errors='replace')
                summary = generate_summary(text_data, max_tokens, AZURE_OPENAI_MODEL, user_name, session_id)
                return {file_name: summary}
            except Exception as e:
                error_msg = f"Error processing text file: {str(e)}"
                log_context('error', user_name, session_id, error_msg)
                return {file_name: f"Error processing text file: {str(e)}"}
        else:
            error_msg = f"Unsupported file format: {file_extension}"
            log_context('warning', user_name, session_id, error_msg)
            return {file_name: error_msg}
    except Exception as e:
        error_msg = f"Error processing document {file_name}: {str(e)}"
        log_context('error', user_name, session_id, error_msg)
        return {file_name: f"Error processing document: {str(e)}"}


def batch_process_documents(document_urls: List[str], max_tokens: int = MAX_SUMMARY_TOKENS, 
                          user_name: str = "system", session_id: str = "system") -> Dict[str, str]:
    """
    Process multiple documents and return their summaries.
    
    Args:
        document_urls (List[str]): List of document URLs
        max_tokens (int, optional): Maximum length of each summary. Defaults to MAX_SUMMARY_TOKENS.
        user_name (str, optional): User identifier for logging. Defaults to "system".
        session_id (str, optional): Session identifier for logging. Defaults to "default".
        
    Returns:
        Dict[str, str]: Dictionary with filenames as keys and summaries as values
    """
    results = {}
    
    for doc_url in document_urls:
        try:
            file_name = doc_url.split('/')[-1]
            log_context('info', user_name, session_id, f"Processing document: {file_name}")
            doc_summary = process_document(doc_url, max_tokens, user_name, session_id)
            results.update(doc_summary)
        except Exception as e:
            file_name = doc_url.split('/')[-1]
            error_msg = f"Failed to process {file_name}: {str(e)}"
            log_context('error', user_name, session_id, error_msg)
            results[file_name] = f"Failed to process: {str(e)}"
    
    return results


def generate_file_summary(user_name: str, session_id: str, files: List[Tuple[str, bytes]], max_tokens: int = MAX_SUMMARY_TOKENS) -> Dict[str, str]:
    """
    End-to-end function that saves files to blob storage and generates summaries.
    
    Args:
        user_name (str): Name of the root container (usually user identifier)
        session_id (str): Unique identifier for the virtual directory
        files (list): List of tuples containing (filename, file_contents)
        max_tokens (int, optional): Maximum length of each summary. Defaults to MAX_SUMMARY_TOKENS.
        
    Returns:
        Dict[str, str]: Dictionary with filenames as keys and summaries as values
        
    Raises:
        Exception: For any errors during file processing or summary generation
    """
    try:
        log_context('info', user_name, session_id, "Starting file summary generation")
        
        # Step 1: Save all files to blob storage
        document_urls = save_files_to_blob_storage(user_name, session_id, files)
        log_context('info', user_name, session_id, f"Successfully saved {len(document_urls)} files to blob storage")
        
        if not document_urls:
            log_context('warning', user_name, session_id, "No valid files were saved to blob storage")
            return {"error": "No valid files were uploaded"}
            
        # Step 2: Process all documents and generate summaries
        summaries = batch_process_documents(document_urls, max_tokens, user_name, session_id)
        log_context('info', user_name, session_id, f"Successfully generated summaries for {len(summaries)} documents")
        
        return summaries
        
    except Exception as e:
        error_msg = f"Error in file summary generation: {str(e)}"
        log_context('error', user_name, session_id, error_msg)
        return {"error": f"Failed to generate summaries: {str(e)}"}
